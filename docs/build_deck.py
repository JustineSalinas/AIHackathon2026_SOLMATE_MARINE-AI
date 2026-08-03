"""Build the Final Technical Pitch Deck (.pptx) from `docs/DECK.md`.

    python docs/build_deck.py

The brief requires the deck in PPT format. Keeping the *content* in Markdown and
generating the PPTX means the deck cannot silently drift from the repository the
way a hand-maintained slide file would -- which matters here more than usual,
because DECK.md's whole premise is that every claim on a slide is one the source
code will confirm. Edit DECK.md; re-run this; never edit the .pptx by hand.

What it does with each part of a slide section:

    **On slide:**      -> the slide body. Bullets, nesting and Markdown tables.
    **Speaker notes:** -> the notes pane.
    **Serves:**        -> appended to the notes as the rubric mapping.
    > Build status...  -> appended to the notes, flagged. These blockquotes say
                          "delete before pitch", so they must never reach a
                          projected slide -- but losing them would lose the
                          measured demo beats they carry, so they land in the
                          notes where only the presenter sees them.

The dark palette is the product's own (`apps/console/src/index.css`): a bridge
instrument at 05:40, not a corporate template. It also survives a projector in a
lit room better than dark-on-white body text does.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
DECK_MD = DOCS / "DECK.md"
OUT = DOCS / "Marine-AI-Pitch-Deck.pptx"

# --- palette -----------------------------------------------------------------
# Same five roles as the console. Colour means state there; here it means
# emphasis, and it is used sparingly for the same reason.
BASE = RGBColor(0x02, 0x06, 0x17)  # page
SURFACE = RGBColor(0x0B, 0x12, 0x20)  # panels, table body
HAIRLINE = RGBColor(0x1E, 0x29, 0x3B)
INK = RGBColor(0xF8, 0xFA, 0xFC)
INK_MUTED = RGBColor(0xB6, 0xC2, 0xD4)  # clears 4.5:1 on BASE
ADVISORY = RGBColor(0xF9, 0x73, 0x16)  # the one accent

SANS = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)  # 16:9
MARGIN = Inches(0.72)


# --- markdown ----------------------------------------------------------------


def split_sections(text: str) -> list[tuple[str, str]]:
    """`## ` headings to (title, body). Front matter before the first one is dropped."""
    parts = re.split(r"^## ", text, flags=re.M)[1:]
    out = []
    for part in parts:
        head, _, body = part.partition("\n")
        out.append((head.strip(), body))
    return out


def strip_links(s: str) -> str:
    """`[text](url)` -> `text`. A projected slide cannot be clicked."""
    return re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", s)


def inline_runs(paragraph, text: str, size: Pt, color=INK, bold_all=False) -> None:
    """Render `**bold**` and `` `code` `` as real runs rather than literal marks."""
    for chunk in re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", strip_links(text)):
        if not chunk:
            continue
        run = paragraph.add_run()
        run.font.size = size
        run.font.name = SANS
        run.font.color.rgb = color
        run.font.bold = bold_all
        if chunk.startswith("**") and chunk.endswith("**"):
            run.text = chunk[2:-2]
            run.font.bold = True
        elif chunk.startswith("`") and chunk.endswith("`"):
            run.text = chunk[1:-1]
            run.font.name = MONO
            run.font.size = Pt(size.pt - 2)
            run.font.color.rgb = INK_MUTED
        else:
            run.text = chunk


def unwrap(body: str) -> list[str]:
    """Re-join bullets that DECK.md wrapped for an 80-column editor.

    Without this, every continuation line becomes its own paragraph -- rendered
    as a phantom sub-bullet -- and any `**bold**` span that happened to straddle
    the wrap is split across two paragraphs, so the asterisks print literally on
    the slide. Both were visible on the ethics slide before this existed.

    A line continues the previous bullet when it is indented, carries no bullet
    marker of its own, and is not part of a table.
    """
    out: list[str] = []
    for raw in body.splitlines():
        line = raw.rstrip()
        stripped = line.strip()
        is_bullet = bool(re.match(r"^\s*[-*]\s+", line))
        is_table = stripped.startswith("|")
        continues = (
            out
            and out[-1].strip()
            and not is_bullet
            and not is_table
            and stripped
            and line.startswith(" ")
            and not out[-1].strip().startswith("|")
        )
        if continues:
            out[-1] = out[-1].rstrip() + " " + stripped
        else:
            out.append(line)
    return out


def parse_table(lines: list[str], i: int) -> tuple[list[list[str]] | None, int]:
    """A GFM table starting at `lines[i]`, or (None, i)."""
    if i >= len(lines) or not lines[i].strip().startswith("|"):
        return None, i
    if i + 1 >= len(lines) or not re.match(r"^\s*\|[\s:|-]+\|\s*$", lines[i + 1]):
        return None, i
    rows = []
    j = i
    while j < len(lines) and lines[j].strip().startswith("|"):
        if not re.match(r"^\s*\|[\s:|-]+\|\s*$", lines[j]):
            cells = [c.strip() for c in lines[j].strip().strip("|").split("|")]
            rows.append(cells)
        j += 1
    return rows, j


# --- chrome ------------------------------------------------------------------


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BASE
    return slide


# ---------------------------------------------------------------- diagrams --
#
# Drawn as native PowerPoint shapes rather than exported images, for the same
# reason the banca mark is: a picture would be a binary asset to keep in sync
# with a deck whose whole premise is that it cannot drift from the repository.
# Shapes also stay crisp on a projector at any resolution, and a presenter can
# nudge a box five minutes before walking on stage.


def _box(slide, x, y, w, h, title, sub=None, *, accent=False, muted=False):
    """One labelled node. `accent` marks the module a slide is really about."""
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = SURFACE
    shp.line.color.rgb = ADVISORY if accent else HAIRLINE
    shp.line.width = Pt(1.5 if accent else 1)
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.12

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(12.5)
    r.font.bold = True
    r.font.name = SANS
    r.font.color.rgb = INK_MUTED if muted else INK

    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(9.5)
        r2.font.name = MONO
        r2.font.color.rgb = ADVISORY if accent else INK_MUTED
    return shp


def _arrow(slide, x, y, w, h=Pt(9)):
    from pptx.enum.shapes import MSO_SHAPE

    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = HAIRLINE
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def _caption(slide, x, y, w, text, *, size=Pt(11), color=INK_MUTED, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.42))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    inline_runs(p, text, size, color)
    return tb


def _label(slide, x, y, w, text):
    """A column heading over a group of nodes."""
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.name = MONO
    r.font.color.rgb = ADVISORY
    return tb


def diagram_architecture(slide, top: Emu) -> Emu:
    """Sense -> decide -> show, with the authority boundary made visible."""
    y = Emu(int(top))
    colw = Inches(3.15)
    x1, x2, x3 = MARGIN, Inches(5.0), Inches(9.5)

    _label(slide, x1, y, colw, "sense")
    _label(slide, x2, y, Inches(3.6), "decide  ·  deterministic")
    _label(slide, x3, y, colw, "show")

    ny = Emu(int(y + Inches(0.36)))
    bh = Inches(0.62)
    gap = Inches(0.18)

    sensors = [
        ("Engine loom", "7 streams · 1 Hz"),
        ("GNSS", "position · speed"),
        ("Metocean", "wind · wave · current"),
    ]
    for i, (t, s) in enumerate(sensors):
        _box(slide, x1, Emu(int(ny + (bh + gap) * i)), colw, bh, t, s)

    modules = [
        ("Speed Optimization", "physics + XGBoost wear", True),
        ("Route Optimization", "same fuel model per leg", True),
        ("Predictive Maintenance", "z-score + PCA autoencoder", True),
        ("Safety cut-offs", "rules only · imports no model", False),
    ]
    for i, (t, s, acc) in enumerate(modules):
        _box(slide, x2, Emu(int(ny + (bh + gap) * i)), Inches(3.6), bh, t, s, accent=acc)

    _box(slide, x3, Emu(int(ny + Inches(0.4))), colw, Inches(1.4),
         "One bridge display", "console · captain view")
    _box(slide, x3, Emu(int(ny + Inches(2.1))), colw, bh,
         "Advisory phrasing", "Claude · guarded", muted=True)

    ay = Emu(int(ny + Inches(1.05)))
    _arrow(slide, Emu(int(x1 + colw + Inches(0.12))), ay, Inches(0.55))
    _arrow(slide, Emu(int(x2 + Inches(3.6) + Inches(0.12))), ay, Inches(0.55))

    cy = Emu(int(ny + (bh + gap) * 4 + Inches(0.12)))
    _caption(
        slide, MARGIN, cy, Emu(int(W - 2 * MARGIN)),
        "**The models predict and detect. They never decide.** The optimiser picks the RPM, "
        "the planner picks the track, a rule table trips the alarms — and a test fails if a "
        "model is ever imported into the safety path.",
    )
    return Emu(int(cy + Inches(0.5)))


def diagram_sensor_bridge(slide, top: Emu) -> Emu:
    """Why Problem 1 and Problem 2 are the same measurement."""
    y = Emu(int(top))
    mid = Inches(4.9)
    bw = Inches(3.5)

    _box(slide, mid, y, bw, Inches(0.78),
         "Exhaust gas temperature", "one sensor, over baseline", accent=True)

    ry = Emu(int(y + Inches(1.35)))
    _box(slide, MARGIN, ry, Inches(3.3), Inches(0.78),
         "Robust z-score + PCA", "learns healthy correlations")
    _box(slide, Inches(9.35), ry, Inches(3.3), Inches(0.78),
         "XGBoost wear model", "1,326 degradation states")

    oy = Emu(int(ry + Inches(1.15)))
    _box(slide, MARGIN, oy, Inches(3.3), Inches(0.78),
         "\"Coolant is drifting\"", "names the stream, never a date")
    _box(slide, Inches(9.35), oy, Inches(3.3), Inches(0.78),
         "+ litres / hour", "the wear priced in fuel")

    ay = Emu(int(y + Inches(0.26)))
    left = slide.shapes.add_shape(13, Inches(4.05), ay, Inches(0.7), Pt(9))  # LEFT_ARROW
    left.fill.solid()
    left.fill.fore_color.rgb = HAIRLINE
    left.line.fill.background()
    left.shadow.inherit = False
    _arrow(slide, Emu(int(mid + bw + Inches(0.1))), ay, Inches(0.7))

    _caption(
        slide, MARGIN, Emu(int(oy + Inches(0.95))), Emu(int(W - 2 * MARGIN)),
        "A worn engine runs hotter at the same load. The **detector** reads that as an anomaly; "
        "the **fuel model** prices it in litres per hour. So \"your engine is degrading\" and "
        "\"it is costing you money\" are not two claims — they are one measurement, read twice.",
    )
    return Emu(int(oy + Inches(1.5)))


DIAGRAMS = {
    "architecture": diagram_architecture,
    "sensor-bridge": diagram_sensor_bridge,
}


def banca_mark(slide, left: Emu, top: Emu, size: Emu) -> None:
    """The outrigger-banca mark, same geometry as the console header.

    Two float spindles and a hull, drawn as shapes rather than an image so the
    deck has no binary asset to keep in sync. Booms are dropped at this size for
    the same reason the favicon drops them.
    """
    from pptx.enum.shapes import MSO_SHAPE

    unit = size / 24.0  # the SVG viewBox is 24 wide

    for cx in (-9.6, 9.6):
        f = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(int(left + (cx - 1.5 + 12) * unit)),
            Emu(int(top + (-8.8 + 13) * unit)),
            Emu(int(3.0 * unit)),
            Emu(int(16.8 * unit)),
        )
        f.fill.solid()
        f.fill.fore_color.rgb = ADVISORY
        f.line.fill.background()
        f.shadow.inherit = False

    hull = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Emu(int(left + (-4.6 + 12) * unit)),
        Emu(int(top + (-12.2 + 13) * unit)),
        Emu(int(9.2 * unit)),
        Emu(int(22.2 * unit)),
    )
    hull.fill.solid()
    hull.fill.fore_color.rgb = ADVISORY
    hull.line.fill.background()
    hull.shadow.inherit = False


def add_footer(slide, prs, index: int, total: int) -> None:
    rule = slide.shapes.add_shape(
        1,  # rectangle
        MARGIN,
        Emu(int(H - Inches(0.72))),
        Emu(int(W - 2 * MARGIN)),
        Pt(1),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = HAIRLINE
    rule.line.fill.background()
    rule.shadow.inherit = False

    box = slide.shapes.add_textbox(
        MARGIN, Emu(int(H - Inches(0.62))), Emu(int(W - 2 * MARGIN)), Inches(0.32)
    )
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Marine-AI  ·  Team SOLMATE  ·  ADVISORY ONLY — CAPTAIN COMMANDS"
    r.font.size = Pt(10)
    r.font.name = SANS
    r.font.color.rgb = INK_MUTED

    num = slide.shapes.add_textbox(
        Emu(int(W - MARGIN - Inches(1.2))), Emu(int(H - Inches(0.62))), Inches(1.2), Inches(0.32)
    )
    np_ = num.text_frame.paragraphs[0]
    np_.alignment = PP_ALIGN.RIGHT
    nr = np_.add_run()
    nr.text = f"{index}/{total}"
    nr.font.size = Pt(10)
    nr.font.name = MONO
    nr.font.color.rgb = INK_MUTED


def set_notes(slide, blocks: list[str]) -> None:
    text = "\n\n".join(b.strip() for b in blocks if b and b.strip())
    if text:
        slide.notes_slide.notes_text_frame.text = text


# --- slide builders ----------------------------------------------------------


def title_slide(prs, title: str, body_lines: list[str], notes: list[str]) -> None:
    slide = add_slide(prs)
    banca_mark(slide, MARGIN, Inches(1.9), Inches(1.5))

    box = slide.shapes.add_textbox(
        Emu(int(MARGIN + Inches(2.0))), Inches(2.0), Emu(int(W - MARGIN - Inches(2.6))), Inches(3.4)
    )
    tf = box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Marine-AI"
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.name = SANS
    r.font.color.rgb = INK
    p.space_after = Pt(10)

    for line in body_lines[1:]:
        line = line.lstrip("-* ").strip()
        if not line:
            continue
        para = tf.add_paragraph()
        para.space_after = Pt(8)
        inline_runs(para, line, Pt(17), INK_MUTED)

    set_notes(slide, notes)


def content_slide(prs, title: str, body: str, notes: list[str], index: int, total: int) -> None:
    slide = add_slide(prs)

    tbox = slide.shapes.add_textbox(MARGIN, Inches(0.52), Emu(int(W - 2 * MARGIN)), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = strip_links(title)
    tr.font.size = Pt(30)
    tr.font.bold = True
    tr.font.name = SANS
    tr.font.color.rgb = INK

    accent = slide.shapes.add_shape(1, MARGIN, Inches(1.44), Inches(1.1), Pt(3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ADVISORY
    accent.line.fill.background()
    accent.shadow.inherit = False

    top = Inches(1.75)
    lines = unwrap(body)
    # Dense slides step down one size rather than overflowing off the bottom
    # edge, which is the failure nobody notices until it is projected.
    dense = sum(len(ln) for ln in lines) > 780
    lead, sub = (Pt(16), Pt(13.5)) if dense else (Pt(18), Pt(15))
    i = 0
    body_box = None
    tf = None

    while i < len(lines):
        raw = lines[i]
        table, j = parse_table(lines, i)
        if table:
            add_table(slide, table, top)
            top = Emu(int(top + Inches(0.42) * len(table) + Inches(0.3)))
            tf = None
            i = j
            continue

        line = raw.strip()
        if not line:
            i += 1
            continue

        # `<!-- diagram:name -->` on its own line renders a native-shape diagram
        # in place. Kept as an HTML comment so DECK.md still reads as prose
        # anywhere Markdown is rendered -- the deck is the document's second
        # output, not its only one.
        dm = re.fullmatch(r"<!--\s*diagram:([a-z0-9-]+)\s*-->", line)
        if dm:
            draw = DIAGRAMS.get(dm.group(1))
            if draw is None:
                raise SystemExit(f"unknown diagram '{dm.group(1)}' in slide: {title}")
            top = draw(slide, top)
            tf = None
            i += 1
            continue

        if tf is None:
            body_box = slide.shapes.add_textbox(
                MARGIN, top, Emu(int(W - 2 * MARGIN)), Emu(int(H - top - Inches(0.9)))
            )
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            first = True

        indent = (len(raw) - len(raw.lstrip(" "))) // 2
        text = re.sub(r"^[-*]\s+", "", line)
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.level = min(indent, 4)
        para.space_after = Pt(9)
        size = lead if indent == 0 else sub
        if line.startswith(("-", "*")):
            bullet = para.add_run()
            bullet.text = "— " if indent == 0 else "· "
            bullet.font.size = size
            bullet.font.name = SANS
            bullet.font.color.rgb = ADVISORY if indent == 0 else INK_MUTED
        inline_runs(para, text, size, INK if indent == 0 else INK_MUTED)
        i += 1

    add_footer(slide, prs, index, total)
    set_notes(slide, notes)


def add_table(slide, rows: list[list[str]], top: Emu) -> None:
    n_rows, n_cols = len(rows), len(rows[0])
    width = Emu(int(W - 2 * MARGIN))
    height = Emu(int(Inches(0.42) * n_rows))
    shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, width, height)
    table = shape.table
    table.first_row = True

    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row[:n_cols]):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE if r else HAIRLINE
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            inline_runs(p, cell_text, Pt(13), INK if r == 0 else INK_MUTED, bold_all=(r == 0))


# --- main --------------------------------------------------------------------


def main() -> int:
    md = DECK_MD.read_text(encoding="utf-8")
    sections = split_sections(md)

    slides = []
    for title, body in sections:
        if title.startswith(("Judging weights", "Pre-submission checklist")):
            continue  # internal, never projected

        # Blockquoted build-status notes carry measured demo beats but say
        # "delete before pitch" -- pull them out of the body, into the notes.
        build_notes = [
            re.sub(r"^>\s?", "", m, flags=re.M).strip()
            for m in re.findall(r"(?:^>.*\n?)+", body, flags=re.M)
        ]
        body = re.sub(r"(?:^>.*\n?)+", "", body, flags=re.M)

        on_slide = ""
        m = re.search(r"\*\*On slide[^:]*:\*\*\n(.*?)(?=\n\*\*Speaker notes|\Z)", body, re.S)
        if m:
            on_slide = m.group(1)
        elif not re.search(r"\*\*Speaker notes", body):
            on_slide = body  # e.g. the appendix list

        notes = []
        m = re.search(r"\*\*Speaker notes:\*\*(.*?)(?=\n\*\*Serves|\Z)", body, re.S)
        if m:
            notes.append(m.group(1).strip())
        m = re.search(r"\*\*Serves:\*\*(.*?)(?=\n\n|\Z)", body, re.S)
        if m:
            notes.append("SERVES (rubric): " + " ".join(m.group(1).split()))
        for bn in build_notes:
            notes.append("BUILD NOTE — resolve before pitching, do not read aloud:\n" + bn)

        slides.append((title, on_slide, notes))

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    total = len(slides) - 1  # the title slide is not numbered
    for n, (title, body, notes) in enumerate(slides):
        clean = re.sub(r"^Slide \d+\s*—\s*", "", title)
        if n == 0:
            title_slide(prs, clean, [ln.strip() for ln in body.splitlines() if ln.strip()], notes)
        else:
            content_slide(prs, clean, body, notes, n, total)

    prs.save(OUT)
    print(f"wrote {OUT.relative_to(DOCS.parent)}  ({len(slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
